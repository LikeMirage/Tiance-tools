from __future__ import annotations

from datetime import datetime
import os
from pathlib import Path
import shutil
from typing import Any

from pptx import Presentation

from tiance_runtime import run_tool
from ppt_editing import apply_operations, inspect_deck, normalize_slide_partnames, select_slides_by_sequence
from ppt_elements import add_slide, merged_theme, set_page_size
from ppt_inventory import quality_warnings, shape_inventory, static_text_risk_report
from ppt_package import inspect_package, prune_unreferenced_parts, validate_saved_pptx


class ToolError(Exception):
    def __init__(self, code: str, message: str, details: dict[str, Any] | None = None) -> None:
        super().__init__(message)
        self.code = code
        self.message = message
        self.details = details or {}


def ok(summary: str, data: dict[str, Any], warnings: list[str] | None = None) -> dict[str, Any]:
    return {"ok": True, "summary": summary, "data": data, "warnings": warnings or []}


def fail(
    code: str,
    message: str,
    details: dict[str, Any] | None = None,
    warnings: list[str] | None = None,
) -> dict[str, Any]:
    return {
        "ok": False,
        "error": f"{code}: {message}",
        "error_info": {"code": code, "message": message, "details": details or {}},
        "warnings": warnings or [],
    }


def workspace_root() -> Path:
    raw = os.environ.get("TIANCE_WORKSPACE_ROOT") or os.environ.get("WORKSPACE_ROOT") or os.getcwd()
    return Path(raw).expanduser().resolve(strict=False)


def read_bool(value: Any, default: bool = False) -> bool:
    if isinstance(value, bool):
        return value
    if isinstance(value, str):
        lowered = value.strip().lower()
        if lowered in {"true", "1", "yes", "y", "on"}:
            return True
        if lowered in {"false", "0", "no", "n", "off"}:
            return False
    return default


def resolve_input_path(value: Any, root: Path) -> Path:
    return resolve_existing_pptx(value, root, "input_path")


def resolve_template_path(value: Any, root: Path) -> Path:
    return resolve_existing_pptx(value, root, "template_path")


def resolve_existing_pptx(value: Any, root: Path, field_name: str) -> Path:
    if not isinstance(value, str) or not value.strip():
        raise ToolError("INVALID_ARGUMENT", f"{field_name} 必须是非空字符串。")
    path = Path(value.strip()).expanduser()
    if not path.is_absolute():
        path = root / path
    resolved = path.resolve(strict=False)
    ensure_inside_workspace(resolved, root, field_name)
    if not resolved.is_file():
        raise ToolError("INPUT_NOT_FOUND", "输入 PPTX 文件不存在。", {field_name: str(resolved)})
    if resolved.suffix.lower() != ".pptx":
        raise ToolError("UNSUPPORTED_FORMAT", "当前只支持 .pptx 文件。", {field_name: str(resolved)})
    return resolved


def resolve_output_path(value: Any, root: Path) -> Path:
    if not isinstance(value, str) or not value.strip():
        raise ToolError("INVALID_ARGUMENT", "output_path 必须是非空字符串。")
    path = Path(value.strip()).expanduser()
    if not path.is_absolute():
        path = root / path
    resolved = path.resolve(strict=False)
    ensure_inside_workspace(resolved, root, "output_path")
    if resolved.suffix.lower() != ".pptx":
        resolved = resolved.with_suffix(".pptx")
    return resolved


def ensure_inside_workspace(path: Path, root: Path, field_name: str) -> None:
    try:
        path.relative_to(root)
    except ValueError as exc:
        raise ToolError(
            "PATH_OUTSIDE_WORKSPACE",
            f"{field_name} 必须位于工作区内。",
            {field_name: str(path), "workspace_root": str(root)},
        ) from exc


def ensure_can_write(output_path: Path, *, overwrite: bool) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    if output_path.exists() and not overwrite:
        raise ToolError("OUTPUT_EXISTS", "输出文件已存在。", {"output_path": str(output_path)})


def apply_core_properties(prs: Presentation, value: Any) -> None:
    if not isinstance(value, dict):
        return
    mapping = {
        "title": "title",
        "subject": "subject",
        "author": "author",
        "category": "category",
        "comments": "comments",
        "keywords": "keywords",
    }
    for key, attr in mapping.items():
        item = value.get(key)
        if isinstance(item, str):
            setattr(prs.core_properties, attr, item)


def create_deck(payload: dict[str, Any], root: Path) -> dict[str, Any]:
    output_path = resolve_output_path(payload.get("output_path"), root)
    overwrite = read_bool(payload.get("overwrite"), False)
    ensure_can_write(output_path, overwrite=overwrite)
    template_path = None
    if payload.get("template_path") is not None:
        template_path = resolve_template_path(payload.get("template_path"), root)
        if template_path == output_path:
            raise ToolError("INVALID_ARGUMENT", "template_path 和 output_path 不能是同一个文件。")

    slides = payload.get("slides")
    has_slides = isinstance(slides, list) and len(slides) > 0
    if template_path is None and not has_slides:
        raise ToolError("INVALID_ARGUMENT", "create 操作必须提供非空 slides。")
    if slides is not None and not isinstance(slides, list):
        raise ToolError("INVALID_ARGUMENT", "slides 必须是数组。")

    theme = merged_theme(payload.get("theme"))
    prs = Presentation(str(template_path)) if template_path else Presentation()
    if template_path is None or payload.get("page") is not None:
        set_page_size(prs, payload.get("page"))
    selected_template_slides: list[int] | None = None
    if template_path and payload.get("template_slide_indexes") is not None:
        raw_indexes = payload.get("template_slide_indexes")
        if not isinstance(raw_indexes, list):
            raise ToolError("INVALID_ARGUMENT", "template_slide_indexes 必须是数组。")
        selected_template_slides = select_slides_by_sequence(prs, raw_indexes)
    if template_path:
        normalize_slide_partnames(prs)
    apply_core_properties(prs, payload.get("properties"))
    for slide_spec in slides or []:
        if not isinstance(slide_spec, dict):
            raise ToolError("INVALID_ARGUMENT", "slides 的每一项都必须是对象。")
        add_slide(prs, slide_spec, theme, root)
    prs.save(output_path)
    warnings: list[str] = []
    finalize_data = finalize_saved_pptx(output_path, warnings)
    return ok(
        f"PPTX 创建完成：{output_path.name}。",
        {
            "action": "create",
            "output_path": str(output_path),
            "template_path": str(template_path) if template_path else "",
            "template_slide_indexes": selected_template_slides,
            "slide_count": finalize_data["validation"]["slide_count"],
            "overwrite": overwrite,
            **finalize_data,
        },
        warnings,
    )


def inspect_ppt(payload: dict[str, Any], root: Path) -> dict[str, Any]:
    input_path = resolve_input_path(payload.get("input_path"), root)
    prs = Presentation(str(input_path))
    inspect_options = payload.get("inspect") if isinstance(payload.get("inspect"), dict) else {}
    max_chars = inspect_options.get("max_text_chars_per_slide", 2000)
    if not isinstance(max_chars, int):
        max_chars = 2000
    data = inspect_deck(
        prs,
        include_text_runs=read_bool(inspect_options.get("include_text_runs"), False),
        max_text_chars_per_slide=max(200, min(max_chars, 10000)),
    )
    if read_bool(inspect_options.get("include_package_metadata"), True):
        data["package"] = inspect_package(input_path)
    if read_bool(inspect_options.get("include_shape_inventory"), False):
        max_shape_chars = inspect_options.get("max_text_chars_per_shape", 1200)
        if not isinstance(max_shape_chars, int):
            max_shape_chars = 1200
        data["shape_inventory"] = shape_inventory(prs, max_text_chars_per_shape=max(200, min(max_shape_chars, 10000)))
    if read_bool(inspect_options.get("include_static_risks"), True):
        data["static_text_risks"] = static_text_risk_report(prs)
    data.update({"action": "inspect", "input_path": str(input_path)})
    return ok(f"PPTX 检查完成：{input_path.name}。", data)


def edit_deck(payload: dict[str, Any], root: Path) -> dict[str, Any]:
    input_path = resolve_input_path(payload.get("input_path"), root)
    output_path = resolve_output_path(payload.get("output_path"), root)
    overwrite = read_bool(payload.get("overwrite"), False)
    backup = read_bool(payload.get("backup"), True)
    dry_run = read_bool(payload.get("dry_run"), False)
    operations = payload.get("operations")
    if not isinstance(operations, list) or not operations:
        raise ToolError("INVALID_ARGUMENT", "edit 操作必须提供非空 operations。")
    if not dry_run:
        ensure_can_write(output_path, overwrite=overwrite)

    prs = Presentation(str(input_path))
    before_count = len(prs.slides)
    theme = merged_theme(payload.get("theme"))
    operation_summaries = apply_operations(prs, operations, theme, root)
    risk_report = static_text_risk_report(prs)
    backup_path = ""
    if dry_run:
        return ok(
            "PPTX 编辑预演完成，未写入文件。",
            {
                "action": "edit",
                "dry_run": True,
                "input_path": str(input_path),
                "output_path": str(output_path),
                "slide_count_before": before_count,
                "slide_count_after": len(prs.slides),
                "operations": operation_summaries,
                "static_text_risks": risk_report,
            },
            quality_warnings(risk_report),
        )

    if input_path == output_path and output_path.exists() and backup:
        backup_path = create_backup(output_path)
    prs.save(output_path)
    warnings: list[str] = []
    finalize_data = finalize_saved_pptx(output_path, warnings)
    return ok(
        f"PPTX 编辑完成：{output_path.name}。",
        {
            "action": "edit",
            "input_path": str(input_path),
            "output_path": str(output_path),
            "backup_path": backup_path,
            "slide_count_before": before_count,
            "slide_count_after": finalize_data["validation"]["slide_count"],
            "operations": operation_summaries,
            "overwrite": overwrite,
            **finalize_data,
        },
        warnings,
    )


def finalize_saved_pptx(output_path: Path, warnings: list[str]) -> dict[str, Any]:
    cleanup_data: dict[str, Any] = {"removed_count": 0, "removed_parts": []}
    try:
        cleanup_data = prune_unreferenced_parts(output_path)
    except Exception as exc:  # noqa: BLE001
        warnings.append(f"PPTX 包清理未完成：{str(exc) or type(exc).__name__}")

    validation = validate_saved_pptx(output_path)
    if not validation.get("ok"):
        raise ToolError("OUTPUT_VALIDATION_FAILED", "输出 PPTX 保存后未通过完整性检查。", validation)

    prs = Presentation(str(output_path))
    risk_report = static_text_risk_report(prs)
    warnings.extend(quality_warnings(risk_report))
    return {
        "cleanup": cleanup_data,
        "validation": validation,
        "static_text_risks": risk_report,
    }


def create_backup(path: Path) -> str:
    stamp = datetime.now().strftime("%Y%m%d%H%M%S")
    backup_path = path.with_suffix(path.suffix + f".{stamp}.bak")
    shutil.copy2(path, backup_path)
    return str(backup_path)


def run(payload: dict[str, Any]) -> dict[str, Any]:
    try:
        action = str(payload.get("action") or "").strip().lower()
        root = workspace_root()
        if action == "create":
            return create_deck(payload, root)
        if action == "inspect":
            return inspect_ppt(payload, root)
        if action == "edit":
            return edit_deck(payload, root)
        raise ToolError("INVALID_ARGUMENT", "action 必须是 create、inspect 或 edit。", {"action": action})
    except ToolError as exc:
        return fail(exc.code, exc.message, exc.details)
    except Exception as exc:
        return fail("PPT_TOOL_FAILED", str(exc) or type(exc).__name__)


if __name__ == "__main__":
    run_tool(run)
