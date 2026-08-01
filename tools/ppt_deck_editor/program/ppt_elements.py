from __future__ import annotations

from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


DEFAULT_THEME = {
    "primary_color": "1F4E79",
    "accent_color": "00A6A6",
    "background_color": "F7F9FC",
    "text_color": "172033",
    "muted_color": "64748B",
    "font_family": "Microsoft YaHei",
}


def merged_theme(value: Any) -> dict[str, Any]:
    theme = dict(DEFAULT_THEME)
    if isinstance(value, dict):
        theme.update({key: item for key, item in value.items() if item is not None})
    return theme


def set_page_size(prs: Presentation, page: Any) -> None:
    size = "wide"
    width = None
    height = None
    if isinstance(page, dict):
        size = str(page.get("size") or "wide")
        width = page.get("width")
        height = page.get("height")
    if size == "standard":
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        return
    if size == "custom" and isinstance(width, (int, float)) and isinstance(height, (int, float)):
        prs.slide_width = Inches(width)
        prs.slide_height = Inches(height)
        return
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def add_slide(prs: Presentation, spec: dict[str, Any], theme: dict[str, Any], root: Path) -> Any:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, spec.get("background"), theme)
    layout = str(spec.get("layout") or "title_content")
    if layout == "title":
        add_title_layout(slide, spec, theme)
    elif layout == "section":
        add_section_layout(slide, spec, theme)
    elif layout == "two_column":
        add_title(slide, spec.get("title"), theme)
    elif layout == "image_focus":
        add_image_focus_layout(slide, spec, theme, root)
    elif layout == "table":
        add_title(slide, spec.get("title"), theme)
        if isinstance(spec.get("table"), dict):
            add_table(slide, spec["table"], theme)
    elif layout != "blank":
        add_title_content_layout(slide, spec, theme)
    for element in spec.get("elements") or []:
        if isinstance(element, dict):
            add_element(slide, element, theme, root)
    return slide


def apply_background(slide: Any, spec: Any, theme: dict[str, Any]) -> None:
    color = None
    if isinstance(spec, dict):
        color = spec.get("color")
    rgb = parse_color(color or theme["background_color"])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_title_layout(slide: Any, spec: dict[str, Any], theme: dict[str, Any]) -> None:
    add_text(
        slide,
        {
            "text": spec.get("title") or "",
            "x": 0.9,
            "y": 2.25,
            "w": 11.5,
            "h": 0.8,
            "style": {"font_size": 34, "bold": True, "align": "center", "color": theme["text_color"]},
        },
        theme,
    )
    if isinstance(spec.get("subtitle"), str):
        add_text(
            slide,
            {
                "text": spec["subtitle"],
                "x": 1.2,
                "y": 3.15,
                "w": 10.9,
                "h": 0.55,
                "style": {"font_size": 16, "align": "center", "color": theme["muted_color"]},
            },
            theme,
        )
    add_shape(
        slide,
        {
            "shape": "rect",
            "x": 5.2,
            "y": 4.0,
            "w": 2.9,
            "h": 0.08,
            "style": {"fill_color": theme["accent_color"], "line_color": theme["accent_color"]},
        },
        theme,
    )


def add_section_layout(slide: Any, spec: dict[str, Any], theme: dict[str, Any]) -> None:
    add_shape(
        slide,
        {
            "shape": "rect",
            "x": 0,
            "y": 0,
            "w": 13.333,
            "h": 7.5,
            "style": {"fill_color": theme["primary_color"], "line_color": theme["primary_color"]},
        },
        theme,
    )
    add_text(
        slide,
        {
            "text": spec.get("title") or "",
            "x": 1.0,
            "y": 2.7,
            "w": 11.2,
            "h": 0.8,
            "style": {"font_size": 32, "bold": True, "color": "FFFFFF", "align": "center"},
        },
        theme,
    )
    if isinstance(spec.get("subtitle"), str):
        add_text(
            slide,
            {
                "text": spec["subtitle"],
                "x": 1.4,
                "y": 3.65,
                "w": 10.5,
                "h": 0.5,
                "style": {"font_size": 16, "color": "DDEBFF", "align": "center"},
            },
            theme,
        )


def add_title_content_layout(slide: Any, spec: dict[str, Any], theme: dict[str, Any]) -> None:
    add_title(slide, spec.get("title"), theme)
    if isinstance(spec.get("body"), str):
        add_text(
            slide,
            {"text": spec["body"], "x": 0.9, "y": 1.55, "w": 11.6, "h": 4.8, "style": {"font_size": 18}},
            theme,
        )
    if isinstance(spec.get("bullets"), list):
        add_bullets(slide, {"items": spec["bullets"], "x": 1.0, "y": 1.65, "w": 11.0, "h": 4.9}, theme)


def add_image_focus_layout(slide: Any, spec: dict[str, Any], theme: dict[str, Any], root: Path) -> None:
    add_title(slide, spec.get("title"), theme)
    if isinstance(spec.get("image_path"), str):
        add_image(
            slide,
            {
                "image_path": spec["image_path"],
                "x": 0.9,
                "y": 1.45,
                "w": 7.1,
                "h": 5.3,
                "sizing": spec.get("image_sizing") or spec.get("sizing") or "contain",
            },
            theme,
            root,
        )
    if isinstance(spec.get("bullets"), list):
        add_bullets(slide, {"items": spec["bullets"], "x": 8.3, "y": 1.7, "w": 4.2, "h": 4.5}, theme)


def add_title(slide: Any, title: Any, theme: dict[str, Any]) -> None:
    if not isinstance(title, str) or not title.strip():
        return
    add_text(
        slide,
        {
            "text": title,
            "x": 0.7,
            "y": 0.45,
            "w": 11.9,
            "h": 0.55,
            "style": {"font_size": 25, "bold": True, "color": theme["text_color"]},
        },
        theme,
    )
    add_shape(
        slide,
        {
            "shape": "rect",
            "x": 0.7,
            "y": 1.12,
            "w": 1.25,
            "h": 0.055,
            "style": {"fill_color": theme["accent_color"], "line_color": theme["accent_color"]},
        },
        theme,
    )


def add_element(slide: Any, element: dict[str, Any], theme: dict[str, Any], root: Path) -> None:
    element_type = str(element.get("type") or "").lower()
    if element_type == "text":
        add_text(slide, element, theme)
    elif element_type == "bullets":
        add_bullets(slide, element, theme)
    elif element_type == "image":
        add_image(slide, element, theme, root)
    elif element_type == "table":
        add_table(slide, element, theme)
    elif element_type == "shape":
        add_shape(slide, element, theme)
    elif element_type == "line":
        add_line(slide, element, theme)
    else:
        raise ValueError(f"不支持的 PPT 元素类型：{element_type}")


def add_text(slide: Any, spec: dict[str, Any], theme: dict[str, Any]) -> Any:
    box = slide.shapes.add_textbox(_in(spec, "x", 0.8), _in(spec, "y", 1.2), _in(spec, "w", 11.6), _in(spec, "h", 1.0))
    text_frame = box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    style = spec.get("style") if isinstance(spec.get("style"), dict) else {}
    paragraph.alignment = paragraph_alignment(style.get("align"))
    run = paragraph.add_run()
    run.text = str(spec.get("text") or "")
    apply_font(run.font, style, theme)
    return box


def add_bullets(slide: Any, spec: dict[str, Any], theme: dict[str, Any]) -> Any:
    box = slide.shapes.add_textbox(_in(spec, "x", 0.9), _in(spec, "y", 1.5), _in(spec, "w", 11.0), _in(spec, "h", 4.8))
    text_frame = box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    style = spec.get("style") if isinstance(spec.get("style"), dict) else {}
    items = spec.get("items") or []
    for index, item in enumerate(items):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        if isinstance(item, dict):
            text = str(item.get("text") or "")
            level = int(item.get("level") or 0)
        else:
            text = str(item)
            level = 0
        paragraph.level = max(0, min(level, 4))
        paragraph.space_after = Pt(8)
        run = paragraph.add_run()
        run.text = text
        apply_font(run.font, style, theme)
    return box


def add_image(slide: Any, spec: dict[str, Any], theme: dict[str, Any], root: Path) -> Any:
    raw_path = spec.get("image_path")
    if not isinstance(raw_path, str) or not raw_path.strip():
        raise ValueError("image 元素缺少 image_path。")
    image_path = Path(raw_path).expanduser()
    if not image_path.is_absolute():
        image_path = root / image_path
    image_path = image_path.resolve(strict=False)
    try:
        image_path.relative_to(root)
    except ValueError as exc:
        raise ValueError(f"图片路径不在工作区内：{image_path}") from exc
    if not image_path.is_file():
        raise ValueError(f"图片文件不存在：{image_path}")
    sizing = image_sizing(spec)
    x = float_value(spec, "x", 0.9)
    y = float_value(spec, "y", 1.3)
    width = float_value(spec, "w", 6.0)
    height = float_value(spec, "h", 4.0)
    if sizing == "stretch":
        return slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), width=Inches(width), height=Inches(height))

    image_width, image_height = image_dimensions(image_path)
    image_ratio = image_width / image_height
    box_ratio = width / height
    if sizing == "cover":
        picture = slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), width=Inches(width), height=Inches(height))
        if image_ratio > box_ratio:
            crop = (1 - (box_ratio / image_ratio)) / 2
            picture.crop_left = crop
            picture.crop_right = crop
        elif image_ratio < box_ratio:
            crop = (1 - (image_ratio / box_ratio)) / 2
            picture.crop_top = crop
            picture.crop_bottom = crop
        return picture

    if image_ratio > box_ratio:
        display_width = width
        display_height = width / image_ratio
        display_x = x
        display_y = y + (height - display_height) / 2
    else:
        display_height = height
        display_width = height * image_ratio
        display_x = x + (width - display_width) / 2
        display_y = y
    return slide.shapes.add_picture(
        str(image_path),
        Inches(display_x),
        Inches(display_y),
        width=Inches(display_width),
        height=Inches(display_height),
    )


def add_table(slide: Any, spec: dict[str, Any], theme: dict[str, Any]) -> Any:
    rows = spec.get("rows")
    if not isinstance(rows, list) or not rows:
        raise ValueError("table 元素缺少 rows。")
    col_count = max(len(row) for row in rows if isinstance(row, list))
    if col_count <= 0:
        raise ValueError("table.rows 不能为空。")
    shape = slide.shapes.add_table(len(rows), col_count, _in(spec, "x", 0.9), _in(spec, "y", 1.5), _in(spec, "w", 11.5), _in(spec, "h", 3.5))
    table = shape.table
    for row_index, row in enumerate(rows):
        if not isinstance(row, list):
            continue
        for col_index in range(col_count):
            cell = table.cell(row_index, col_index)
            cell.text = "" if col_index >= len(row) or row[col_index] is None else str(row[col_index])
            style = spec.get("header_style") if row_index == 0 else spec.get("style")
            format_cell_text(cell, style if isinstance(style, dict) else {}, theme)
            fill_color = None
            if row_index == 0:
                fill_color = style.get("fill_color") if isinstance(style, dict) else None
                fill_color = fill_color or theme["primary_color"]
            elif isinstance(style, dict):
                fill_color = style.get("fill_color")
            if fill_color:
                cell.fill.solid()
                cell.fill.fore_color.rgb = parse_color(fill_color)
    return shape


def add_shape(slide: Any, spec: dict[str, Any], theme: dict[str, Any]) -> Any:
    shape_type = shape_type_from_name(str(spec.get("shape") or "rect"))
    shape = slide.shapes.add_shape(shape_type, _in(spec, "x", 0.8), _in(spec, "y", 1.2), _in(spec, "w", 2.0), _in(spec, "h", 0.8))
    style = spec.get("style") if isinstance(spec.get("style"), dict) else {}
    fill_color = style.get("fill_color") or spec.get("fill_color")
    line_color = style.get("line_color") or spec.get("line_color") or fill_color
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = parse_color(fill_color)
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = parse_color(line_color)
    if isinstance(style.get("line_width"), (int, float)):
        shape.line.width = Pt(style["line_width"])
    if isinstance(spec.get("text"), str):
        shape.text = spec["text"]
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                apply_font(run.font, style, theme)
    return shape


def add_line(slide: Any, spec: dict[str, Any], theme: dict[str, Any]) -> Any:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        _in(spec, "x", 0.8),
        _in(spec, "y", 1.2),
        _in(spec, "x2", float(spec.get("x", 0.8)) + float(spec.get("w", 2.0))),
        _in(spec, "y2", float(spec.get("y", 1.2)) + float(spec.get("h", 0.0))),
    )
    style = spec.get("style") if isinstance(spec.get("style"), dict) else {}
    line.line.color.rgb = parse_color(style.get("line_color") or theme["accent_color"])
    if isinstance(style.get("line_width"), (int, float)):
        line.line.width = Pt(style["line_width"])
    return line


def format_cell_text(cell: Any, style: dict[str, Any], theme: dict[str, Any]) -> None:
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = paragraph_alignment(style.get("align"))
        for run in paragraph.runs:
            apply_font(run.font, style, theme)


def apply_font(font: Any, style: dict[str, Any], theme: dict[str, Any]) -> None:
    font.name = str(style.get("font_family") or theme["font_family"])
    font.size = Pt(float(style.get("font_size") or 16))
    font.bold = bool(style.get("bold")) if "bold" in style else None
    font.italic = bool(style.get("italic")) if "italic" in style else None
    font.color.rgb = parse_color(style.get("color") or theme["text_color"])


def paragraph_alignment(value: Any) -> Any:
    mapping = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    return mapping.get(str(value or "left").lower(), PP_ALIGN.LEFT)


def shape_type_from_name(value: str) -> Any:
    mapping = {
        "rect": MSO_SHAPE.RECTANGLE,
        "rectangle": MSO_SHAPE.RECTANGLE,
        "round_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "diamond": MSO_SHAPE.DIAMOND,
        "cloud": MSO_SHAPE.CLOUD,
        "hexagon": MSO_SHAPE.HEXAGON,
    }
    return mapping.get(value.lower(), MSO_SHAPE.RECTANGLE)


def parse_color(value: Any) -> RGBColor:
    raw = str(value or "000000").strip().lstrip("#")
    if len(raw) == 8:
        raw = raw[-6:]
    if len(raw) != 6 or any(char not in "0123456789abcdefABCDEF" for char in raw):
        raw = "000000"
    return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def image_sizing(spec: dict[str, Any]) -> str:
    style = spec.get("style") if isinstance(spec.get("style"), dict) else {}
    value = str(spec.get("sizing") or style.get("sizing") or "contain").strip().lower()
    if value not in {"contain", "cover", "stretch"}:
        return "contain"
    return value


def image_dimensions(path: Path) -> tuple[float, float]:
    with Image.open(path) as image:
        width, height = image.size
    if width <= 0 or height <= 0:
        raise ValueError(f"图片尺寸无效：{path}")
    return float(width), float(height)


def float_value(spec: dict[str, Any], key: str, default: float) -> float:
    value = spec.get(key, default)
    try:
        return float(value)
    except (TypeError, ValueError):
        return default


def _in(spec: dict[str, Any], key: str, default: float) -> Any:
    value = spec.get(key, default)
    try:
        return Inches(float(value))
    except (TypeError, ValueError):
        return Inches(default)
